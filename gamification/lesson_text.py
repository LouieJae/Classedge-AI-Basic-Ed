"""[Classedge LMS] Extract plain text from uploaded lesson files."""
import os
import re

MAX_CHARS = 20_000
MIN_USABLE_CHARS = 200


class UnsupportedFileType(Exception):
    pass


class EmptyContent(Exception):
    pass


def _clean(text: str) -> str:
    text = re.sub(r"\s+", " ", text).strip()
    return text[:MAX_CHARS]


def _from_pdf(file_obj) -> str:
    from pypdf import PdfReader
    reader = PdfReader(file_obj)
    return "\n".join(p.extract_text() or "" for p in reader.pages)


def _from_docx(file_obj) -> str:
    from docx import Document
    doc = Document(file_obj)
    return "\n".join(p.text for p in doc.paragraphs)


def _from_pptx(file_obj) -> str:
    from pptx import Presentation
    prs = Presentation(file_obj)
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                out.append(shape.text)
    return "\n".join(out)


def _from_txt(file_obj) -> str:
    raw = file_obj.read()
    if isinstance(raw, bytes):
        raw = raw.decode("utf-8", errors="ignore")
    return raw


_HANDLERS = {".pdf": _from_pdf, ".docx": _from_docx, ".pptx": _from_pptx, ".txt": _from_txt}


def extract_text(file_obj) -> str:
    name = getattr(file_obj, "name", "") or ""
    ext = os.path.splitext(name)[1].lower()
    handler = _HANDLERS.get(ext)
    if handler is None:
        raise UnsupportedFileType(f"Unsupported file type: {ext or 'unknown'}")
    raw = handler(file_obj)
    cleaned = _clean(raw)
    if len(cleaned) < MIN_USABLE_CHARS:
        raise EmptyContent("Lesson file appears empty or unreadable.")
    return cleaned
